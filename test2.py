from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
bullet_slide_layout = prs.slide_layouts[0]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title

left = top = width = height = Inches(1.5)
print(left)
txBox = slide.shapes.add_textbox(left, top + 60, width, height)
tf = txBox.text_frame
tf.text = "that's my text" 
# body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'

# tf = body_shape.text_frame
# tf.text = 'Find the bullet slide layout'
# tf.level = -1

# p = tf.add_paragraph()
# p.text = 'Use _TextFrame.text for first bullet'
# p.level = 1

# p = tf.add_paragraph()
# p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
# p.level = 2

prs.save('test.pptx')