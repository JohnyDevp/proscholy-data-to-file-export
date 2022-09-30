import json
import pptx

from pptx import Presentation
from props import bcolors
from aspose import slides

class Presentation:
    def __init__(self, json_content):
        self.json_content = json_content
    
    def __createPPTSlideShow(self):
        # create new presentation and pass it to ppt-creating function
        pptxPres = pptx.Presentation()

        # go through each slide and create 
        
        # create first slide with exact layout
        layout = pptxPres.slide_layouts[0] 
        slide = pptxPres.slides.add_slide(layout)
        shapes = slide.shapes
        shapes.title.text = self.json_content['slides'][0]['title']
        shapes.placeholders[1].text = self.json_content['slides'][0]['sub-title']
        
        # iterate through all next slides
        for sl in self.json_content['slides'][1:]:
            # TODO set a font for the text
            layout = pptxPres.slide_layouts[0] 
            slide = pptxPres.slides.add_slide(layout)
            shapes = slide.shapes
            shapes.title.text = sl['content']

        # return it
        return pptxPres


    def exportPPT(self):
        pptxPres = self.__createPPTSlideShow()
        
        # save the presentation whith its name
        pptxPres.save(self.json_content['file']['name']+'.ppt')
        
        # return the name
        return self.json_content['file']['name']+'.ppt'
    
    def exportPDF(self):
        pres_name = self.exportPPT()
        
        new_name = str(pres_name).replace('.ppt', '.pdf')
        # load the created presentation for converting to pdf
        load_pres = slides.Presentation(pres_name)
        load_pres.save(new_name, slides.export.SaveFormat.PDF)

        # return the new name (so, replaced ppt extension for pdf)
        return new_name